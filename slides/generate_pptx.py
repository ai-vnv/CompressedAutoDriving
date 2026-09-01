"""Editable PPTX version of the internship-summary deck (mirrors main.tex).

Theme: AI V&V Lab hero palette — emerald gradient background, white bold
headings, gold accents, mint body, Poppins. Poppins must be installed on the
viewing machine (TTFs are in slides/fonts/) or PowerPoint substitutes a
default sans.  Gold spans are marked with **...** in the bullet strings.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
FIG = HERE / "figures"

EMERALD_DEEP = RGBColor(0x07, 0x33, 0x2A)
EMERALD_LIGHT = RGBColor(0x15, 0x60, 0x4B)
GOLD = RGBColor(0xC9, 0xA4, 0x5C)
MINT = RGBColor(0xD9, 0xE7, 0xE1)
MINT_DIM = RGBColor(0xA8, 0xC4, 0xBA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)
FONT = "Poppins"
FONT_BOLD = "Poppins SemiBold"


def add_background(slide):
    shp = slide.shapes.add_shape(1, 0, 0, SW, SH)
    shp.line.fill.background()
    fill = shp.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = EMERALD_LIGHT
    stops[0].position = 0.0
    stops[1].color.rgb = EMERALD_DEEP
    stops[1].position = 1.0
    try:
        fill.gradient_angle = 65.0
    except Exception:
        pass
    shp.shadow.inherit = False


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def set_runs(par, segments, size, color, bold=False, font=FONT):
    for text, is_gold in segments:
        run = par.add_run()
        run.text = text
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold or is_gold
        f.color.rgb = GOLD if is_gold else color


def parse(text):
    out, gold, buf = [], False, ""
    i = 0
    while i < len(text):
        if text.startswith("**", i):
            if buf:
                out.append((buf, gold))
            buf, gold = "", not gold
            i += 2
        else:
            buf += text[i]
            i += 1
    if buf:
        out.append((buf, gold))
    return out


def kicker(slide, text, y=Inches(0.55), align=PP_ALIGN.LEFT, x=Inches(0.75),
           w=Inches(11.8)):
    tf = textbox(slide, x, y, w, Inches(0.35))
    p = tf.paragraphs[0]
    p.alignment = align
    set_runs(p, [(" ".join(text.upper()), False)], 10, GOLD, bold=True)


def title(slide, text, y=Inches(0.95), size=28):
    tf = textbox(slide, Inches(0.75), y, Inches(11.8), Inches(0.9))
    p = tf.paragraphs[0]
    set_runs(p, parse(text), size, WHITE, bold=True, font=FONT_BOLD)


def bullets(slide, items, y, w=Inches(11.8), x=Inches(0.75), size=16,
            gap=Pt(8)):
    tf = textbox(slide, x, y, w, SH - y - Inches(0.6))
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = gap
        set_runs(p, [("•  ", True)] + parse(item), size, MINT)


def footnote(slide, text, y, size=13, color=None):
    tf = textbox(slide, Inches(0.75), y, Inches(11.8), Inches(0.8))
    p = tf.paragraphs[0]
    set_runs(p, parse(text), size, color or MINT_DIM)


def picture_card(slide, name, x, y, w=None, h=None, pad=Inches(0.08)):
    from PIL import Image

    img = Image.open(FIG / name)
    ar = img.width / img.height
    if w is None:
        w = Emu(int(h * ar))
    if h is None:
        h = Emu(int(w / ar))
    card = slide.shapes.add_shape(1, x - pad, y - pad,
                                  Emu(int(w) + 2 * pad), Emu(int(h) + 2 * pad))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()
    card.shadow.inherit = False
    slide.shapes.add_picture(str(FIG / name), x, y, width=w, height=h)
    return w, h


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def centered(slide, text, y, size, color, bold=False, font=FONT):
    tf = textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_runs(p, parse(text), size, color, bold=bold, font=font)


prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    return s


# ==== ACT I ====================================================================

# 1 · title
s = new_slide()
kicker(s, "AI Verification & Validation Lab — Internship Summary",
       y=Inches(1.55), align=PP_ALIGN.CENTER, x=Inches(0.6), w=Inches(12.1))
centered(s, "From **Backend Engineer**", Inches(2.15), 34, WHITE, True, FONT_BOLD)
centered(s, "to Closed-Loop **V&V** Research", Inches(2.85), 34, WHITE, True,
         FONT_BOLD)
centered(s, "Learning, building, and evaluating decision-making policies "
            "under uncertainty", Inches(3.7), 14, MINT_DIM)
centered(s, "Ahmad Alfan Alfian Irfan", Inches(4.7), 15, MINT)
centered(s, "AI VnV Research Intern", Inches(5.2), 11, MINT_DIM)
note(s, "One sentence: this is the story of the whole internship — from "
        "backend engineering to a submitted V&V paper.")

# 2 · starting point
s = new_slide()
kicker(s, "The starting point")
title(s, "A backend engineer, curious about uncertainty")
bullets(s, [
    "Background: backend engineering — not computer-science research",
    "The curiosity: how do agents **decide under uncertainty**? (POMDPs)",
    "Month one, honestly: unfamiliar ground — the theory, the tooling, the papers",
    "The result was **bad performance** — and a decision to fix it properly",
], Inches(2.1))
footnote(s, "The gap was fundamentals, not effort.", Inches(5.6))
note(s, "Be honest about the first month — it sets up the whole arc. ~1 min.")

# 3 · the reset
s = new_slide()
kicker(s, "The reset")
title(s, "Back to fundamentals, deliberately")
bullets(s, [
    "Diagnosed the gap: decision-making theory and its ecosystem",
    "**Coursera** — decision-making course, from MDP basics up",
    "**JuliaAcademy** — Decision Making Under Uncertainty with POMDPs.jl",
    "Rule adopted: every concept gets **rebuilt in code** before moving on",
], Inches(2.1))
footnote(s, "Learn → build → verify became the working loop for everything "
            "after.", Inches(5.6))
note(s, "~45 s.")

# 4 · learning by building
s = new_slide()
kicker(s, "Learning by building")
title(s, "Pluto notebooks as the learning medium")
bullets(s, [
    "Hands-on **Pluto notebooks**: MDP and POMDP from first principles",
    "Interactive **solver playground** — explore the JuliaPOMDP solver "
    "ecosystem side by side",
    "Worked through the **Gallery of POMDPs.jl** examples",
    "The notebooks double as a learning medium for the next intern",
], Inches(2.1))
footnote(s, "**github.com/PannnTastic/Decision-Making-Under-Uncertainty**",
         Inches(5.5), size=14)
note(s, "~1 min. Mention the notebooks are interactive — sliders, live "
        "solvers.")

# 5 · project 1: DuckieMDP
s = new_slide()
kicker(s, "First project — DuckieMDP")
title(s, "Certified explanations for MDP driving policies")
picture_card(s, "mdp_explain.png", Inches(3.4), Inches(1.85), h=Inches(2.9))
bullets(s, [
    "MDP formulation of Duckietown driving; explanations in three pillars: "
    "**why this action** · **what if another** · **is it consistent**",
    "Counterfactual rollouts share exogenous noise — outcome claims settled "
    "**exactly**, across Q-learning, SARSA, SAC, TD3",
], Inches(5.0), size=13, gap=Pt(5))
footnote(s, "**github.com/PannnTastic/DuckieMDP**  · arXiv preprint",
         Inches(6.45), size=13)
note(s, "First real project: fully observed MDP, certified explanation "
        "framework. ~1.25 min.")

# 6 · side project: Duckietown.jl
s = new_slide()
kicker(s, "Side project — Duckietown.jl")
title(s, "The simulator, reborn in Julia")
picture_card(s, "duckietown_jl.png", Inches(0.75), Inches(2.3), w=Inches(6.3))
footnote(s, "DORA solver completing a lap, drawn by the package's native "
            "renderer", Inches(5.1), size=10)
bullets(s, [
    "Native **POMDPs.jl** reimplementation of the DuckieMDP environment",
    "Validated **decision by decision** against Python — exact NumPy RNG "
    "streams, seeded episodes reproduce **bit for bit**",
    "Native renderer; DORASolvers.jl case study drives a full lap",
], Inches(2.4), x=Inches(7.45), w=Inches(5.3), size=13, gap=Pt(8))
tf = textbox(s, Inches(7.45), Inches(5.5), Inches(5.3), Inches(0.4))
set_runs(tf.paragraphs[0], parse("**github.com/ai-vnv/Duckietown.jl**"),
         13, GOLD)
note(s, "The RNG-parity point is the punchline — the Julia port reproduces "
        "the Python episodes bit for bit. ~1 min.")

# ==== ACT II ===================================================================

# 7 · divider
s = new_slide()
kicker(s, "Main project — submitted to IEEE OJ-ITS", y=Inches(1.9),
       align=PP_ALIGN.CENTER, x=Inches(0.6), w=Inches(12.1))
centered(s, "A **Closed-Loop** Evaluation of", Inches(2.5), 28, WHITE, True,
         FONT_BOLD)
centered(s, "Capability **Loss and Recovery**", Inches(3.1), 28,
         WHITE, True, FONT_BOLD)
centered(s, "in Compressed Driving Policies", Inches(3.7), 28, WHITE, True,
         FONT_BOLD)
centered(s, "From learning POMDPs to using one as a research instrument.",
         Inches(4.6), 13, MINT_DIM)
note(s, "Transition into the main project. ~30 s.")

# 8 · problem
s = new_slide()
kicker(s, "The problem")
title(s, "A policy can pass the benchmark and fail to stop")
bullets(s, [
    "Learned driving policies run on embedded computers",
    "Standard practice: prune, distill, quantize before deployment",
    "Progress measured by size, latency, aggregate accuracy",
    "Aggregate scores hide **which behaviors** compression removes",
], Inches(2.1))
footnote(s, "The question is not only how small the network can get — but "
            "which learned capabilities survive each stage.", Inches(5.6))
note(s, "~45 s.")

# 9 · key idea
s = new_slide()
kicker(s, "Key idea")
title(s, "Closed loop, stage by stage — reversed")
bullets(s, [
    "A policy acts on its own observations: small action drift **compounds** "
    "through the feedback loop",
    "Static test-set scores cannot see this",
    "Scenario-based assessment, reversed: **fix the scenarios and seeds**, "
    "vary the system under test",
    "Every compression stage is a new configuration to re-certify; acceptance "
    "criteria fixed **before** any result",
], Inches(2.1))
footnote(s, "Four questions: where does driving break · is it recoverable · "
            "does the data decide · does recovery survive lower precision.",
         Inches(5.7))
note(s, "The V&V framing. ~1 min.")

# 10 · task
s = new_slide()
kicker(s, "The driving task")
title(s, "Five curricula, shared seeds, rising difficulty")
picture_card(s, "fig_task.png", Inches(0.85), Inches(2.2), w=Inches(11.6))
bullets(s, [
    "C0–C1 lane following · C2 crossing pedestrian · C3 stop sign · C4 both",
    "The stop curricula demand interrupting the drive, then re-establishing it",
], Inches(4.6), size=14, gap=Pt(5))
note(s, "~45 s.")

# 11 · system under test
s = new_slide()
kicker(s, "System under test")
title(s, "A belief-state visuomotor policy (POMDP + PPO)")
picture_card(s, "fig_pipeline.png", Inches(0.75), Inches(2.15), w=Inches(7.1))
bullets(s, [
    "Camera → **YOLO11n** (objects) + **MobileNetV3** (lane pose)",
    "EKFs fuse detections into a **29-d belief** vector",
    "Belief-PPO actor: **73,986** parameters (FP32)",
    "Only the actor is compressed; perception stays fixed",
], Inches(2.25), x=Inches(8.25), w=Inches(4.5), size=13, gap=Pt(7))
note(s, "~1 min.")

# 12 · how the pipeline works
s = new_slide()
kicker(s, "How the pipeline works")
title(s, "Ten configurations, one shared ancestry")
bullets(s, [
    "**A0** original actor (reference)",
    "**A1** structured pruning, hidden width 256 → 64",
    "**A2** + distillation on the hardest curriculum only",
    "**A3** + distillation on balanced C0–C4 data (62,176 states)",
    "**A6** A3 → post-training INT8    **A8** A3 parent → QAT INT8",
    "**A9** A3 → FP16 cast (precision control)",
    "**A4, A5, A7** controls: unpruned INT8, no-distill INT8, historical QAT",
], Inches(2.1), size=14, gap=Pt(6))
footnote(s, "Every arrow is evaluated the same way — no configuration is "
            "trusted on the previous stage's evidence.", Inches(6.1))
note(s, "Story line: A1 → A2/A3 → A6/A8 → A9. ~1 min.")

# 13 · protocol
s = new_slide()
kicker(s, "Evaluation protocol")
title(s, "400 matched episodes, two independent axes")
bullets(s, [
    "**10 models × 5 curricula × 8 seeds** — identical seeds for all",
    "Acceptance checks preregistered: completion, progress, collisions, "
    "stop compliance and restart, lane keeping, clearance",
    "Evaluation backend verified **bit-for-bit** reproducible",
    "Axis 1: task verdicts (does it drive?)",
    "Axis 2: action fidelity vs. the original on identical inputs (does it imitate?)",
], Inches(2.1))
note(s, "~45 s.")

# 14 · findings 1+2
s = new_slide()
kicker(s, "Findings 1 & 2")
title(s, "Pruning breaks driving; rehearsal decides recovery")
bullets(s, [
    "Structured pruning: 73,986 → **6,210** parameters (−91.6%) — "
    "**fails all five** curricula",
    "Same teacher, loss, optimizer, budget — only the rehearsal data differs",
    "Distill on the hardest curriculum only → recovers **C3–C4 only**",
    "Balanced rehearsal across all five → recovers **all five**",
], Inches(2.1))
footnote(s, "What the student rehearses is what the student keeps.",
         Inches(5.6))
note(s, "~1 min.")

# 15 · finding 3
s = new_slide()
kicker(s, "Finding 3")
title(s, "INT8 undoes a completed recovery")
picture_card(s, "fig_matrix.png", Inches(4.15), Inches(1.95), h=Inches(3.9))
bullets(s, [
    "Both INT8 routes fail C3–C4 from one byte-identical checkpoint",
    "Stop-task completion falls from **8/8 to 3/8** under PTQ",
], Inches(6.2), size=14, gap=Pt(4))
note(s, "Walk the matrix top to bottom. ~1 min.")

# 16 · phenotypes
s = new_slide()
kicker(s, "Opposite failure phenotypes")
title(s, "One checkpoint, two ways to fail a stop")
picture_card(s, "fig_stop.png", Inches(0.85), Inches(1.95), h=Inches(5.1))
bullets(s, [
    "**A6 (PTQ)** brakes correctly, parks **0.22 m** before the line — and "
    "never moves again (to the 2700-step horizon)",
    "**A8 (QAT)** drives **through** the stop at low speed while the stop is "
    "still required",
    "Same parent, same task, opposite failures",
], Inches(2.3), x=Inches(6.3), w=Inches(6.3), size=14, gap=Pt(8))
note(s, "Frames reconstructed from telemetry; circles on the curves are the "
        "same steps. ~1 min.")

# 17 · findings 4+5
s = new_slide()
kicker(s, "Findings 4 & 5")
title(s, "Controls and fidelity: the cause is the interaction")
bullets(s, [
    "Same PTQ on the **unpruned** actor → passes all five",
    "**FP16 cast** of the recovered actor → keeps all five",
    "So the failure is **narrow width × INT8**, not precision alone",
    "QAT imitates the original **more closely** than PTQ — and drives worse",
    "The unpruned INT8 control drives everything — and fails the similarity "
    "thresholds",
], Inches(2.1))
footnote(s, "Action-level similarity is not an acceptance test — in either "
            "direction.", Inches(5.9))
note(s, "~1 min.")

# 18 · cost
s = new_slide()
kicker(s, "What compression buys")
title(s, "Actor cost across precisions")
rows = [
    ("", "Param. memory", "Median latency", "All curricula"),
    ("A3 FP32", "24.8 KB", "19.4 µs", "pass"),
    ("A9 FP16", "**12.4 KB**", "24.5 µs (+26%)", "pass"),
    ("A6 INT8", "**7.6 KB**", "**12.9 µs** (1.50×)", "fail C3–C4"),
]
tbl_tf = textbox(s, Inches(1.6), Inches(2.15), Inches(10.2), Inches(2.2))
for i, row in enumerate(rows):
    p = tbl_tf.paragraphs[0] if i == 0 else tbl_tf.add_paragraph()
    p.space_after = Pt(6)
    text = "     ".join(f"{c:<22}" for c in row)
    set_runs(p, parse(text), 15, MINT if i else WHITE, bold=(i == 0))
bullets(s, [
    "FP16 halves memory but is slower here (no native half path on this CPU)",
    "The actor is under **0.1%** of end-to-end step cost — perception dominates",
], Inches(4.7), size=14, gap=Pt(6))
note(s, "~45 s.")

# 19 · takeaways
s = new_slide()
kicker(s, "Takeaways")
title(s, "Judge compressed policies by how they drive")
bullets(s, [
    "Compression is a sequence of **behavioral transitions**: lose, recover, "
    "lose again",
    "Evaluate **every stage** in closed loop — no stage is safe by inheritance",
    "Choose rehearsal data to cover **everything** the policy must keep",
    "A completed recovery is **not** robust to a later precision cut",
    "Action-level similarity is not an acceptance test",
], Inches(2.1))
note(s, "~45 s.")

# 20 · future works
s = new_slide()
kicker(s, "Future works")
title(s, "Where this goes next")
bullets(s, [
    "**Wider settings**: more quantization configurations, more training "
    "realizations, more seeds",
    "**Compress perception**: the other 99.9% of step cost, under the same "
    "acceptance protocol",
    "**Transfer the protocol**: other policy stacks, other simulators, "
    "eventually hardware",
    "**Julia-native line**: Duckietown.jl as the ground for solver-level V&V "
    "experiments",
], Inches(2.1))
note(s, "~1 min.")

# ==== ACT III ==================================================================

# 21 · reflection
s = new_slide()
kicker(s, "What the internship built")
title(s, "Fundamentals first, then honest evaluation")
bullets(s, [
    "From backend engineering to an **RL / POMDP research workflow**",
    "**4 public repos** · **2 papers** · **1 Julia package**",
    "From a bad first month to an **OJ-ITS submission**",
    "The lesson that stuck: measure what matters, **in closed loop**",
], Inches(2.1))
note(s, "Personal close before Q&A. ~45 s.")

# 22 · thanks
s = new_slide()
kicker(s, "Thank you", y=Inches(1.9), align=PP_ALIGN.CENTER, x=Inches(0.6),
       w=Inches(12.1))
centered(s, "Questions?", Inches(2.5), 32, WHITE, True, FONT_BOLD)
centered(s, "Everything in this talk is public:", Inches(3.6), 14, MINT)
centered(s, "**github.com/ai-vnv/CompressedAutoDriving   ·   "
            "github.com/PannnTastic/DuckieMDP**", Inches(4.15), 14, GOLD)
centered(s, "**github.com/PannnTastic/Decision-Making-Under-Uncertainty   ·   "
            "github.com/ai-vnv/Duckietown.jl**", Inches(4.65), 14, GOLD)
centered(s, "Internship summary — AI Verification & Validation Lab",
         Inches(5.6), 11, MINT_DIM)
note(s, "Close and take questions.")

out = HERE / "presentation.pptx"
prs.save(out)
print("saved", out, f"({len(prs.slides._sldIdLst)} slides)")
